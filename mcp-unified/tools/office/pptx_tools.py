"""
PPTX Tools for PowerPoint presentation handling
"""
from pathlib import Path
from typing import Dict, List, Optional, Any

from tools.base import register_tool


@register_tool
def read_pptx(file_path: str) -> Dict:
    """
    Read a PowerPoint file and extract content
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        Dictionary with slides content and metadata
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        slides_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                'slide_number': slide_num,
                'texts': [],
                'shapes_count': len(slide.shapes)
            }
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data['texts'].append(shape.text)
                
                # Extract from tables
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        table_data.append(' | '.join(row_text))
                    if table_data:
                        slide_data['texts'].append('[TABLE]\n' + '\n'.join(table_data))
            
            slides_content.append(slide_data)
        
        # Get presentation metadata
        core_props = prs.core_properties
        
        return {
            'success': True,
            'file_path': file_path,
            'total_slides': len(prs.slides),
            'slides': slides_content,
            'metadata': {
                'title': core_props.title,
                'author': core_props.author,
                'subject': core_props.subject,
                'created': str(core_props.created) if core_props.created else None,
                'modified': str(core_props.modified) if core_props.modified else None,
            }
        }
        
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx not installed. Install with: pip install python-pptx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file_path': file_path
        }


@register_tool
def write_pptx(
    file_path: str,
    slides: List[Dict],
    title: Optional[str] = None
) -> Dict:
    """
    Create a new PowerPoint presentation
    
    Args:
        file_path: Path where the .pptx file will be saved
        slides: List of slide content dictionaries
               Each slide should have:
               - 'title': Slide title
               - 'content': List of bullet points or text
               - 'layout': Optional layout index (default: 1 - Title and Content)
        title: Presentation title (optional)
        
    Returns:
        Success status and file path
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        # Set metadata
        if title:
            prs.core_properties.title = title
        
        for slide_data in slides:
            # Get layout (default to Title and Content)
            layout_idx = slide_data.get('layout', 1)
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 1
            
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if 'title' in slide_data:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_data['title']
            
            # Add content
            if 'content' in slide_data:
                # Find content placeholder
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Body placeholder
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    tf = content_placeholder.text_frame
                    tf.clear()  # Clear existing text
                    
                    for idx, content_item in enumerate(slide_data['content']):
                        if idx == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = content_item
                        p.level = 0
        
        # Save presentation
        prs.save(file_path)
        
        return {
            'success': True,
            'file_path': file_path,
            'message': f'Presentation saved with {len(slides)} slides',
            'total_slides': len(slides)
        }
        
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx not installed. Install with: pip install python-pptx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file_path': file_path
        }


@register_tool
def extract_text_pptx(file_path: str) -> Dict:
    """
    Extract all text from a PowerPoint presentation
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        Dictionary with extracted text
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        all_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f'--- Slide {slide_num} ---']
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            all_text.append('\n'.join(slide_text))
        
        full_text = '\n\n'.join(all_text)
        
        return {
            'success': True,
            'file_path': file_path,
            'total_slides': len(prs.slides),
            'text': full_text,
            'total_characters': len(full_text)
        }
        
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx not installed. Install with: pip install python-pptx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file_path': file_path
        }


@register_tool
def add_slide_pptx(
    file_path: str,
    title: str,
    content: List[str],
    layout: int = 1
) -> Dict:
    """
    Add a slide to an existing PowerPoint presentation
    
    Args:
        file_path: Path to the .pptx file
        title: Slide title
        content: List of bullet points
        layout: Slide layout index (default: 1 - Title and Content)
        
    Returns:
        Success status and details
    """
    try:
        from pptx import Presentation
        
        try:
            prs = Presentation(file_path)
        except:
            prs = Presentation()
        
        # Add slide
        if layout >= len(prs.slide_layouts):
            layout = 1
        
        slide_layout = prs.slide_layouts[layout]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Add content
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Body
                tf = shape.text_frame
                for idx, item in enumerate(content):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                break
        
        prs.save(file_path)
        
        return {
            'success': True,
            'file_path': file_path,
            'slide_added': title,
            'total_slides': len(prs.slides)
        }
        
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx not installed. Install with: pip install python-pptx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file_path': file_path
        }
